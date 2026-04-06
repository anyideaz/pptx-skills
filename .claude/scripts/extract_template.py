"""
extract_template.py — Standalone CLI for PPTX template context extraction.

Usage:
    python extract_template.py <pptx_path> <output_dir>

Output:
    <output_dir>/context.json  — Structured template context
    <output_dir>/images/       — Extracted image files

Exit codes:
    0 — Success
    1 — Missing or invalid arguments
    2 — PPTX file not found
    3 — PPTX parse error
"""
import base64
import colorsys
import hashlib
import io
import json
import os
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


# ── Utility ─────────────────────────────────────────────────────────────────

def _emu_to_inches(emu) -> float:
    try:
        return round(int(emu) / 914400, 3)
    except Exception:
        return 0.0


# ── Theme Color Maps ─────────────────────────────────────────────────────────

_THEME_COLOR_MAP = {
    MSO_THEME_COLOR.ACCENT_1: "accent1",
    MSO_THEME_COLOR.ACCENT_2: "accent2",
    MSO_THEME_COLOR.ACCENT_3: "accent3",
    MSO_THEME_COLOR.ACCENT_4: "accent4",
    MSO_THEME_COLOR.ACCENT_5: "accent5",
    MSO_THEME_COLOR.ACCENT_6: "accent6",
    MSO_THEME_COLOR.DARK_1: "dk1",
    MSO_THEME_COLOR.DARK_2: "dk2",
    MSO_THEME_COLOR.LIGHT_1: "lt1",
    MSO_THEME_COLOR.LIGHT_2: "lt2",
    MSO_THEME_COLOR.HYPERLINK: "hlink",
    MSO_THEME_COLOR.FOLLOWED_HYPERLINK: "folHlink",
}

_DEFAULT_CLR_MAP = {
    "bg1": "lt1",
    "tx1": "dk1",
    "bg2": "lt2",
    "tx2": "dk2",
    "accent1": "accent1",
    "accent2": "accent2",
    "accent3": "accent3",
    "accent4": "accent4",
    "accent5": "accent5",
    "accent6": "accent6",
    "hlink": "hlink",
    "folHlink": "folHlink",
}


def _extract_clr_map(master) -> Dict[str, str]:
    """Extract color-mapping overrides from a slide master's <p:clrMap> element."""
    clr_map = dict(_DEFAULT_CLR_MAP)
    try:
        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        clr_map_elem = master._element.find(f"{{{ns_p}}}clrMap")
        if clr_map_elem is not None:
            for attr in clr_map_elem.attrib:
                clr_map[attr] = clr_map_elem.attrib[attr]
    except Exception:
        pass
    return clr_map


def _build_resolved_theme_colors(
    raw_colors: Dict[str, str],
    clr_map: Dict[str, str],
) -> Dict[str, str]:
    """Resolve logical color names through the master's clrMap.

    PowerPoint documents use logical names (bg1, tx1, bg2…) that map to
    physical names (lt1, dk1, lt2…) via the master's <p:clrMap>.  This
    helper returns a dict where BOTH the logical name AND the physical name
    point to the same resolved hex color, so callers can look up either.
    """
    resolved: Dict[str, str] = {}

    # Start with all raw (physical) colors
    for k, v in raw_colors.items():
        if v:
            resolved[k] = v

    # Add aliases for logical names
    for logical, physical in clr_map.items():
        if physical in raw_colors and raw_colors[physical]:
            resolved[logical] = raw_colors[physical]

    return resolved


# ── Color Resolution ─────────────────────────────────────────────────────────

def _apply_color_modifiers(hex_color: str, color_elem) -> str:
    """Apply OOXML color modifiers (lumMod, lumOff, tint, shade) to a hex color."""
    if not hex_color or len(hex_color) != 6:
        return hex_color
    try:
        r = int(hex_color[0:2], 16) / 255.0
        g = int(hex_color[2:4], 16) / 255.0
        b = int(hex_color[4:6], 16) / 255.0

        h, l, s = colorsys.rgb_to_hls(r, g, b)

        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

        lum_mod = color_elem.find(f"{{{ns_a}}}lumMod")
        lum_off = color_elem.find(f"{{{ns_a}}}lumOff")

        if lum_mod is not None or lum_off is not None:
            mod = int(lum_mod.get("val", "100000")) / 100000.0 if lum_mod is not None else 1.0
            off = int(lum_off.get("val", "0")) / 100000.0 if lum_off is not None else 0.0
            l = min(1.0, max(0.0, l * mod + off))
            r2, g2, b2 = colorsys.hls_to_rgb(h, l, s)
            return "{:02X}{:02X}{:02X}".format(int(r2 * 255), int(g2 * 255), int(b2 * 255))

        tint = color_elem.find(f"{{{ns_a}}}tint")
        shade = color_elem.find(f"{{{ns_a}}}shade")

        if tint is not None:
            t = int(tint.get("val", "100000")) / 100000.0
            r = r + (1.0 - r) * (1.0 - t)
            g = g + (1.0 - g) * (1.0 - t)
            b = b + (1.0 - b) * (1.0 - t)
            return "{:02X}{:02X}{:02X}".format(int(r * 255), int(g * 255), int(b * 255))

        if shade is not None:
            s_val = int(shade.get("val", "100000")) / 100000.0
            r = r * s_val
            g = g * s_val
            b = b * s_val
            return "{:02X}{:02X}{:02X}".format(int(r * 255), int(g * 255), int(b * 255))

    except Exception:
        pass
    return hex_color


def _resolve_color(color_obj, theme_colors: Dict[str, str]) -> Tuple[Optional[str], Optional[str]]:
    """Resolve a python-pptx color object to a hex string and optional theme reference."""
    try:
        if color_obj.type is None:
            return None, None
        if color_obj.type == 1:  # RGB
            return str(color_obj.rgb), None
        if color_obj.type == 2:  # THEME
            theme_ref = _THEME_COLOR_MAP.get(color_obj.theme_color)
            if theme_ref:
                hex_color = theme_colors.get(theme_ref)
                # Apply brightness modifiers if available
                try:
                    elem = color_obj._color
                    if hex_color:
                        hex_color = _apply_color_modifiers(hex_color, elem)
                except Exception:
                    pass
                return hex_color, theme_ref
    except Exception:
        pass
    return None, None


# ── Theme Extraction ─────────────────────────────────────────────────────────

def _parse_theme_xml(theme_xml_bytes: bytes) -> Dict[str, Any]:
    """Parse theme XML bytes into a structured dict with color_scheme and font_scheme."""
    import xml.etree.ElementTree as ET

    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ns_thm = "http://schemas.openxmlformats.org/drawingml/2006/theme"

    result: Dict[str, Any] = {"color_scheme": {}, "font_scheme": {}}

    try:
        root = ET.fromstring(theme_xml_bytes)

        theme_el = root if root.tag == f"{{{ns_thm}}}theme" else root.find(f".//{{{ns_thm}}}theme")
        if theme_el is None:
            theme_el = root

        # Color scheme
        clr_scheme = theme_el.find(f".//{{{ns_a}}}clrScheme")
        if clr_scheme is not None:
            color_map = {}
            for child in clr_scheme:
                color_name = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                color_val = None
                srgb = child.find(f"{{{ns_a}}}srgbClr")
                sys_clr = child.find(f"{{{ns_a}}}sysClr")
                if srgb is not None:
                    color_val = srgb.get("val")
                elif sys_clr is not None:
                    color_val = sys_clr.get("lastClr")
                if color_val:
                    color_map[color_name] = color_val.upper()
            result["color_scheme"] = color_map

        # Font scheme
        font_scheme = theme_el.find(f".//{{{ns_a}}}fontScheme")
        if font_scheme is not None:
            for font_type in ["majorFont", "minorFont"]:
                font_elem = font_scheme.find(f"{{{ns_a}}}{font_type}")
                if font_elem is not None:
                    latin = font_elem.find(f"{{{ns_a}}}latin")
                    if latin is not None:
                        result["font_scheme"][font_type] = latin.get("typeface", "")
    except Exception:
        pass

    return result


def _extract_theme(pptx_bytes: bytes) -> Dict[str, Any]:
    """Extract the primary theme from the PPTX zip."""
    import zipfile
    try:
        with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as z:
            names = z.namelist()
            theme_files = sorted([n for n in names if n.startswith("ppt/theme/theme") and n.endswith(".xml")])
            if theme_files:
                return _parse_theme_xml(z.read(theme_files[0]))
    except Exception:
        pass
    return {"color_scheme": {}, "font_scheme": {}}


def _extract_themes_per_master(pptx_bytes: bytes) -> List[Dict[str, Any]]:
    """Extract theme for each slide master in the PPTX."""
    import zipfile
    import xml.etree.ElementTree as ET

    themes: List[Dict[str, Any]] = []
    try:
        with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as z:
            names = z.namelist()

            # Read master relationship files to find which theme each master uses
            master_rels = sorted([n for n in names if "slideMasters/_rels" in n and n.endswith(".rels")])
            master_theme_files: List[Optional[str]] = []
            for rel_file in master_rels:
                theme_file = None
                try:
                    rel_xml = z.read(rel_file)
                    root = ET.fromstring(rel_xml)
                    for rel in root:
                        target = rel.get("Target", "")
                        if "theme" in target.lower() and target.endswith(".xml"):
                            # Resolve relative path
                            base = "ppt/slideMasters/"
                            if target.startswith("../"):
                                target = "ppt/" + target[3:]
                            elif not target.startswith("ppt/"):
                                target = base + target
                            theme_file = target
                            break
                except Exception:
                    pass
                master_theme_files.append(theme_file)

            for theme_file in master_theme_files:
                if theme_file and theme_file in names:
                    themes.append(_parse_theme_xml(z.read(theme_file)))
                else:
                    themes.append(_extract_theme(pptx_bytes))
    except Exception:
        pass

    if not themes:
        themes.append(_extract_theme(pptx_bytes))

    return themes


# ── Fill & Line Extraction ───────────────────────────────────────────────────

def _extract_fill(shape, part, theme_colors: Dict[str, str]) -> Dict[str, Any]:
    """Extract fill information from a shape."""
    result: Dict[str, Any] = {"type": "none"}
    try:
        fill = shape.fill
        from pptx.enum.dml import MSO_FILL
        if fill.type == MSO_FILL.SOLID:
            color, theme_ref = _resolve_color(fill.fore_color, theme_colors)
            result = {"type": "solid", "color": color}
            if theme_ref:
                result["theme_ref"] = theme_ref
            if color is None:
                try:
                    sp_elem = shape._element
                    solid_fill = sp_elem.find(".//" + qn("a:solidFill"))
                    if solid_fill is not None:
                        scheme_clr = solid_fill.find(qn("a:schemeClr"))
                        if scheme_clr is not None:
                            clr_name = scheme_clr.get("val", "")
                            resolved = theme_colors.get(clr_name)
                            if resolved:
                                resolved = _apply_color_modifiers(resolved, scheme_clr)
                            result["color"] = resolved
                            result["theme_ref"] = clr_name
                except Exception:
                    pass
            try:
                solid_fill = shape._element.find(".//" + qn("a:solidFill"))
                if solid_fill is not None:
                    for clr_elem in solid_fill:
                        alpha_elem = clr_elem.find(qn("a:alpha"))
                        if alpha_elem is not None:
                            alpha_val = int(alpha_elem.get("val", "100000"))
                            transparency = round((100000 - alpha_val) / 1000)
                            if transparency > 0:
                                result["transparency"] = transparency
                            break
            except Exception:
                pass
        elif fill.type == MSO_FILL.GRADIENT:
            stops = []
            try:
                for stop in fill.gradient_stops:
                    c, _ = _resolve_color(stop.color, theme_colors)
                    stops.append({"color": c, "position": round(stop.position, 2)})
            except Exception:
                pass
            result = {"type": "gradient", "stops": stops}
            try:
                grad_fill_el = shape._element.find(".//" + qn("a:gradFill"))
                if grad_fill_el is not None:
                    lin = grad_fill_el.find(qn("a:lin"))
                    if lin is not None:
                        ang = lin.get("ang")
                        if ang:
                            result["angle"] = round(int(ang) / 60000)
                        result["gradient_type"] = "linear"
                    else:
                        path_el = grad_fill_el.find(qn("a:path"))
                        if path_el is not None:
                            result["gradient_type"] = path_el.get("path", "circle")
                        else:
                            result["gradient_type"] = "linear"
            except Exception:
                pass
        elif fill.type == MSO_FILL.PICTURE:
            result = {"type": "image"}
        elif fill.type == MSO_FILL.PATTERNED:
            result = {"type": "pattern"}
        elif fill.type == MSO_FILL.BACKGROUND:
            result = {"type": "inherit"}
    except Exception:
        pass
    return result


def _extract_line_properties(shape, theme_colors: Dict[str, str]) -> Optional[Dict[str, Any]]:
    """Extract line/border properties from a shape's <a:ln> element."""
    try:
        elem = shape._element
        ln = elem.find(".//" + qn("a:ln"))
        if ln is None:
            return None
        if ln.find(qn("a:noFill")) is not None:
            return None

        line_info: Dict[str, Any] = {}
        w = ln.get("w")
        if w:
            line_info["width_pt"] = round(int(w) / 12700, 2)

        prstDash = ln.find(qn("a:prstDash"))
        if prstDash is not None:
            line_info["dash"] = prstDash.get("val", "solid")

        solid_fill = ln.find(qn("a:solidFill"))
        if solid_fill is not None:
            srgb = solid_fill.find(qn("a:srgbClr"))
            if srgb is not None:
                color_val = srgb.get("val")
                if color_val:
                    color_val = _apply_color_modifiers(color_val, srgb)
                line_info["color"] = color_val
            else:
                scheme_clr = solid_fill.find(qn("a:schemeClr"))
                if scheme_clr is not None:
                    clr_name = scheme_clr.get("val", "")
                    resolved = theme_colors.get(clr_name)
                    if resolved:
                        resolved = _apply_color_modifiers(resolved, scheme_clr)
                    line_info["color"] = resolved
                    line_info["color_theme_ref"] = clr_name

        if line_info:
            return line_info
    except Exception:
        pass
    return None


# ── Color from XML ───────────────────────────────────────────────────────────

def _resolve_color_from_xml(color_parent, theme_colors: Dict[str, str]) -> Tuple[Optional[str], Optional[str]]:
    """Resolve color from an XML element's solidFill child."""
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    solid_fill = color_parent.find(f"{{{ns_a}}}solidFill")
    if solid_fill is None:
        return None, None

    srgb = solid_fill.find(f"{{{ns_a}}}srgbClr")
    if srgb is not None:
        color_val = srgb.get("val")
        if color_val:
            color_val = _apply_color_modifiers(color_val, srgb)
        return color_val, None

    scheme_clr = solid_fill.find(f"{{{ns_a}}}schemeClr")
    if scheme_clr is not None:
        clr_name = scheme_clr.get("val", "")
        resolved = theme_colors.get(clr_name)
        if resolved:
            resolved = _apply_color_modifiers(resolved, scheme_clr)
        return resolved, clr_name

    sys_clr = solid_fill.find(f"{{{ns_a}}}sysClr")
    if sys_clr is not None:
        color_val = sys_clr.get("lastClr", sys_clr.get("val", ""))
        if color_val and len(color_val) == 6:
            color_val = _apply_color_modifiers(color_val, sys_clr)
        return color_val or None, None

    hsl_clr = solid_fill.find(f"{{{ns_a}}}hslClr")
    if hsl_clr is not None:
        try:
            h = int(hsl_clr.get("hue", "0")) / 60000.0
            s = int(hsl_clr.get("sat", "0")) / 100000.0
            l = int(hsl_clr.get("lum", "0")) / 100000.0
            r, g, b = colorsys.hls_to_rgb(h / 360.0, l, s)
            color_val = "{:02X}{:02X}{:02X}".format(int(r * 255), int(g * 255), int(b * 255))
            color_val = _apply_color_modifiers(color_val, hsl_clr)
            return color_val, None
        except Exception:
            pass

    scrgb_clr = solid_fill.find(f"{{{ns_a}}}scrgbClr")
    if scrgb_clr is not None:
        try:
            r = int(scrgb_clr.get("r", "0")) / 100000.0
            g = int(scrgb_clr.get("g", "0")) / 100000.0
            b = int(scrgb_clr.get("b", "0")) / 100000.0
            color_val = "{:02X}{:02X}{:02X}".format(
                int(min(1.0, r) * 255), int(min(1.0, g) * 255), int(min(1.0, b) * 255)
            )
            color_val = _apply_color_modifiers(color_val, scrgb_clr)
            return color_val, None
        except Exception:
            pass

    prst_clr = solid_fill.find(f"{{{ns_a}}}prstClr")
    if prst_clr is not None:
        prst_name = prst_clr.get("val", "")
        preset_hex = _PRESET_COLORS.get(prst_name)
        if preset_hex:
            preset_hex = _apply_color_modifiers(preset_hex, prst_clr)
        return preset_hex, None

    return None, None


_PRESET_COLORS = {
    "black": "000000", "white": "FFFFFF", "red": "FF0000", "green": "008000",
    "blue": "0000FF", "yellow": "FFFF00", "cyan": "00FFFF", "magenta": "FF00FF",
    "silver": "C0C0C0", "gray": "808080", "maroon": "800000", "olive": "808000",
    "navy": "000080", "purple": "800080", "teal": "008080", "aqua": "00FFFF",
    "orange": "FFA500", "pink": "FFC0CB", "coral": "FF7F50", "crimson": "DC143C",
    "darkBlue": "00008B", "darkGreen": "006400", "darkRed": "8B0000",
    "lightBlue": "ADD8E6", "lightGreen": "90EE90", "lightGray": "D3D3D3",
    "darkGray": "A9A9A9",
}


# ── Font & Text Extraction ───────────────────────────────────────────────────

def _extract_font_info(run, theme_colors: Optional[Dict[str, str]] = None) -> Dict[str, Any]:
    """Extract font info from a text run."""
    font_info: Dict[str, Any] = {}
    try:
        f = run.font
        if f.name:
            font_info["name"] = f.name
        if f.size is not None:
            font_info["size_pt"] = round(f.size.pt, 1)
        if f.bold is not None:
            font_info["bold"] = f.bold
        if f.italic is not None:
            font_info["italic"] = f.italic

        color_resolved = False
        if theme_colors is not None:
            try:
                rPr = run._r.find(qn("a:rPr"))
                if rPr is not None:
                    color_hex, theme_ref = _resolve_color_from_xml(rPr, theme_colors)
                    if color_hex:
                        font_info["color"] = color_hex
                        if theme_ref:
                            font_info["color_theme_ref"] = theme_ref
                        color_resolved = True
            except Exception:
                pass

        if not color_resolved:
            try:
                if f.color and f.color.rgb:
                    font_info["color"] = str(f.color.rgb)
            except Exception:
                pass
    except Exception:
        pass
    return font_info


def _extract_default_font_from_xml(element, theme_colors: Dict[str, str]) -> Dict[str, Any]:
    """Extract default font properties from XML <a:defRPr> elements."""
    font_info: Dict[str, Any] = {}
    try:
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        def_rpr = None

        lst_style = element.find(f".//{{{ns_a}}}lstStyle")
        if lst_style is not None:
            lvl1 = lst_style.find(f"{{{ns_a}}}lvl1pPr")
            if lvl1 is not None:
                def_rpr = lvl1.find(f"{{{ns_a}}}defRPr")

        if def_rpr is None:
            for p_elem in element.findall(f".//{{{ns_a}}}p"):
                pPr = p_elem.find(f"{{{ns_a}}}pPr")
                if pPr is not None:
                    def_rpr = pPr.find(f"{{{ns_a}}}defRPr")
                    if def_rpr is not None:
                        break

        if def_rpr is None:
            def_rpr = element.find(f".//{{{ns_a}}}endParaRPr")

        if def_rpr is None:
            return font_info

        sz = def_rpr.get("sz")
        if sz:
            font_info["size_pt"] = round(int(sz) / 100, 1)

        b = def_rpr.get("b")
        if b is not None:
            font_info["bold"] = b == "1"

        i_attr = def_rpr.get("i")
        if i_attr is not None:
            font_info["italic"] = i_attr == "1"

        latin = def_rpr.find(f"{{{ns_a}}}latin")
        if latin is not None:
            typeface = latin.get("typeface", "")
            if typeface and not typeface.startswith("+"):
                font_info["name"] = typeface
            elif typeface == "+mj-lt":
                font_info["theme_font"] = "major"
            elif typeface == "+mn-lt":
                font_info["theme_font"] = "minor"

        solid_fill = def_rpr.find(f"{{{ns_a}}}solidFill")
        if solid_fill is not None:
            srgb = solid_fill.find(f"{{{ns_a}}}srgbClr")
            if srgb is not None:
                color_val = srgb.get("val")
                if color_val:
                    color_val = _apply_color_modifiers(color_val, srgb)
                font_info["color"] = color_val
            else:
                scheme_clr = solid_fill.find(f"{{{ns_a}}}schemeClr")
                if scheme_clr is not None:
                    clr_name = scheme_clr.get("val", "")
                    resolved = theme_colors.get(clr_name)
                    if resolved:
                        resolved = _apply_color_modifiers(resolved, scheme_clr)
                    font_info["color"] = resolved
                    font_info["color_theme_ref"] = clr_name
    except Exception:
        pass
    return font_info


def _extract_text_body_properties(element) -> Dict[str, Any]:
    """Extract text body properties (anchor, margins) from <a:bodyPr> element."""
    props: Dict[str, Any] = {}
    try:
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        body_pr = element.find(f".//{{{ns_a}}}bodyPr")
        if body_pr is None:
            return props

        anchor = body_pr.get("anchor")
        if anchor:
            anchor_map = {"t": "top", "ctr": "middle", "b": "bottom"}
            props["vertical_align"] = anchor_map.get(anchor, anchor)

        wrap = body_pr.get("wrap")
        if wrap:
            props["wrap"] = wrap

        for attr, key in [("lIns", "margin_left"), ("rIns", "margin_right"),
                          ("tIns", "margin_top"), ("bIns", "margin_bottom")]:
            val = body_pr.get(attr)
            if val:
                props[key] = round(int(val) / 914400, 2)
    except Exception:
        pass
    return props


def _find_inherited_bullet(ph_shape, level: int) -> Optional[Dict]:
    """Look up inherited bullet for a placeholder shape."""
    _UNSET = object()

    def _check_lst_style(lst_style_elem, lvl: int):
        if lst_style_elem is None:
            return _UNSET
        for try_lvl in (lvl + 1, 1):
            el = lst_style_elem.find(qn(f"a:lvl{try_lvl}pPr"))
            if el is None:
                continue
            if el.find(qn("a:buNone")) is not None:
                return None
            bc = el.find(qn("a:buChar"))
            if bc is not None:
                return {"type": "char", "char": bc.get("char", "•")}
            ba = el.find(qn("a:buAutoNum"))
            if ba is not None:
                return {"type": "auto_num", "scheme": ba.get("type", "")}
        return _UNSET

    try:
        ph_idx = ph_shape.placeholder_format.idx
        layout = ph_shape.part.slide_layout

        for lph in layout.placeholders:
            if lph.placeholder_format.idx != ph_idx:
                continue
            txBody = lph._element.find(qn("p:txBody"))
            if txBody is not None:
                result = _check_lst_style(txBody.find(qn("a:lstStyle")), level)
                if result is not _UNSET:
                    return result
            break

        master = layout.slide_master
        for mph in master.placeholders:
            if mph.placeholder_format.idx != ph_idx:
                continue
            txBody = mph._element.find(qn("p:txBody"))
            if txBody is not None:
                result = _check_lst_style(txBody.find(qn("a:lstStyle")), level)
                if result is not _UNSET:
                    return result
            break
    except Exception:
        pass

    return None


def _extract_paragraph(para, theme_colors: Dict[str, str], ph_shape=None) -> Dict[str, Any]:
    """Extract paragraph-level info."""
    p_info: Dict[str, Any] = {"text": para.text}

    try:
        if para.alignment is not None:
            p_info["alignment"] = str(para.alignment).split(".")[-1].split("(")[0].lower()
    except Exception:
        pass

    try:
        pPr = para._element.find(qn("a:pPr"))
        explicit_bullet_found = False
        level = 0
        if pPr is not None:
            level = int(pPr.get("lvl", "0"))
            buNone = pPr.find(qn("a:buNone"))
            buChar = pPr.find(qn("a:buChar"))
            buAutoNum = pPr.find(qn("a:buAutoNum"))
            if buChar is not None:
                p_info["bullet"] = {"type": "char", "char": buChar.get("char", "")}
                explicit_bullet_found = True
            elif buAutoNum is not None:
                p_info["bullet"] = {"type": "auto_num", "scheme": buAutoNum.get("type", "")}
                explicit_bullet_found = True
            elif buNone is not None:
                p_info["bullet"] = None
                explicit_bullet_found = True

        if not explicit_bullet_found and ph_shape is not None:
            try:
                if ph_shape.placeholder_format is not None:
                    inherited = _find_inherited_bullet(ph_shape, level)
                    if inherited is not None:
                        p_info["bullet"] = inherited
            except Exception:
                pass
    except Exception:
        pass

    if para.runs:
        p_info["font"] = _extract_font_info(para.runs[0], theme_colors)

    if "font" not in p_info or "color" not in p_info.get("font", {}):
        try:
            ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
            fallback_color = None
            fallback_ref = None

            pPr_el = para._element.find(qn("a:pPr"))
            if pPr_el is not None:
                def_rpr = pPr_el.find(qn("a:defRPr"))
                if def_rpr is not None:
                    fallback_color, fallback_ref = _resolve_color_from_xml(def_rpr, theme_colors)

            if fallback_color is None:
                txBody = para._element.getparent()
                if txBody is not None:
                    lst_style = txBody.find(f"{{{ns_a}}}lstStyle")
                    if lst_style is not None:
                        lvl1 = lst_style.find(f"{{{ns_a}}}lvl1pPr")
                        if lvl1 is not None:
                            def_rpr = lvl1.find(f"{{{ns_a}}}defRPr")
                            if def_rpr is not None:
                                fallback_color, fallback_ref = _resolve_color_from_xml(def_rpr, theme_colors)

            if fallback_color:
                font = p_info.setdefault("font", {})
                font["color"] = fallback_color
                if fallback_ref:
                    font["color_theme_ref"] = fallback_ref
        except Exception:
            pass

    try:
        if para.line_spacing is not None:
            p_info["line_spacing_pt"] = round(para.line_spacing.pt, 1) if hasattr(para.line_spacing, "pt") else None
        if para.space_before is not None:
            p_info["space_before_pt"] = round(para.space_before.pt, 1) if hasattr(para.space_before, "pt") else None
        if para.space_after is not None:
            p_info["space_after_pt"] = round(para.space_after.pt, 1) if hasattr(para.space_after, "pt") else None
    except Exception:
        pass

    return p_info


# ── Table Extraction ─────────────────────────────────────────────────────────

def _extract_table(shape, theme_colors: Dict[str, str]) -> Dict[str, Any]:
    """Extract table structure, cell contents, and styling."""
    table = shape.table
    rows_data = []

    for row in table.rows:
        cells = []
        for cell in row.cells:
            cell_info: Dict[str, Any] = {"text": cell.text}
            try:
                cell_fill = cell.fill
                if cell_fill.type is not None:
                    from pptx.enum.dml import MSO_FILL
                    if cell_fill.type == MSO_FILL.SOLID:
                        color, theme_ref = _resolve_color(cell_fill.fore_color, theme_colors)
                        if color:
                            cell_info["fill"] = color
                        if theme_ref:
                            cell_info["fill_theme_ref"] = theme_ref
            except Exception:
                pass
            try:
                if cell.text_frame.paragraphs:
                    para = cell.text_frame.paragraphs[0]
                    if para.runs:
                        fi = _extract_font_info(para.runs[0])
                        if fi:
                            cell_info["font"] = fi
            except Exception:
                pass
            cells.append(cell_info)
        rows_data.append(cells)

    table_info: Dict[str, Any] = {
        "rows": len(table.rows),
        "cols": len(table.columns),
        "data": rows_data,
    }

    try:
        table_info["column_widths_inches"] = [
            _emu_to_inches(col.width) for col in table.columns
        ]
    except Exception:
        pass

    try:
        tbl_pr = shape._element.find(".//" + qn("a:tblPr"))
        if tbl_pr is not None:
            for attr in ("firstRow", "lastRow", "bandRow", "bandCol", "firstCol", "lastCol"):
                val = tbl_pr.get(attr)
                if val is not None:
                    table_info[attr] = val == "1"
    except Exception:
        pass

    return table_info


# ── Image Extraction ─────────────────────────────────────────────────────────

def _get_image_from_rel(part, r_embed: str) -> Optional[Any]:
    """Safely get an image part from a relationship ID."""
    try:
        if r_embed and r_embed in part.rels:
            return part.rels[r_embed].target_part
    except Exception:
        pass
    return None


def _extract_image_from_shape(shape, part, prefix: str, seen_blobs: set,
                               offset_left: int = 0, offset_top: int = 0) -> Optional[Dict]:
    """Try to extract an image from a shape."""
    pos = {
        "left": _emu_to_inches((shape.left or 0) + offset_left),
        "top": _emu_to_inches((shape.top or 0) + offset_top),
        "width": _emu_to_inches(shape.width),
        "height": _emu_to_inches(shape.height),
    }

    try:
        image = shape.image
        blob = image.blob
        blob_hash = hashlib.md5(blob[:512]).hexdigest()
        ext = image.ext if image.ext else "png"
        img_name = f"{prefix}_{shape.name}.{ext}".replace(" ", "_")

        img_entry: Dict[str, Any] = {"name": img_name, **pos}

        if blob_hash not in seen_blobs:
            seen_blobs.add(blob_hash)
            with Image.open(io.BytesIO(blob)) as img:
                w_px, h_px = img.size
            img_entry["media_type"] = f"image/{ext}"
            img_entry["width_px"] = w_px
            img_entry["height_px"] = h_px

        img_entry["base64_data"] = base64.b64encode(blob).decode("utf-8")
        return img_entry
    except Exception:
        pass

    elem = shape._element
    blips = elem.findall(".//" + qn("a:blip"))
    for blip in blips:
        r_embed = blip.get(qn("r:embed"))
        image_part = _get_image_from_rel(part, r_embed)
        if image_part is not None:
            blob = image_part.blob
            ct = image_part.content_type
            if "svg" in ct:
                continue
            blob_hash = hashlib.md5(blob[:512]).hexdigest()
            ext = ct.split("/")[-1].replace("jpeg", "jpg")
            img_name = f"{prefix}_{shape.name}.{ext}".replace(" ", "_")

            img_entry = {"name": img_name, **pos}
            if blob_hash not in seen_blobs:
                seen_blobs.add(blob_hash)
                with Image.open(io.BytesIO(blob)) as img_obj:
                    w_px, h_px = img_obj.size
                img_entry["media_type"] = ct
                img_entry["width_px"] = w_px
                img_entry["height_px"] = h_px

            img_entry["base64_data"] = base64.b64encode(blob).decode("utf-8")
            return img_entry

    SVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
    svg_blips = elem.findall(".//{%s}svgBlip" % SVG_NS)
    for svg_blip in svg_blips:
        r_embed = svg_blip.get(qn("r:embed"))
        image_part = _get_image_from_rel(part, r_embed)
        if image_part is not None:
            blob = image_part.blob
            blob_hash = hashlib.md5(blob[:512]).hexdigest()
            img_name = f"{prefix}_{shape.name}.svg".replace(" ", "_")

            img_entry = {"name": img_name, **pos, "is_svg": True, "media_type": "image/svg+xml"}
            if blob_hash not in seen_blobs:
                seen_blobs.add(blob_hash)

            img_entry["base64_data"] = base64.b64encode(blob).decode("utf-8")
            return img_entry

    return None


# ── Shape Type Name ──────────────────────────────────────────────────────────

def _shape_type_name(shape) -> str:
    """Get a human-readable shape type name."""
    try:
        auto_shape_type = shape.auto_shape_type
        if auto_shape_type is not None:
            return str(auto_shape_type).split(".")[-1].split("(")[0]
    except Exception:
        pass
    try:
        st = shape.shape_type
        if st is not None:
            return str(st).split(".")[-1].split("(")[0]
    except Exception:
        pass
    return "UNKNOWN"


# ── Group Shape Offset ───────────────────────────────────────────────────────

def _get_group_offset_emu(group_shape) -> Tuple[int, int]:
    """Return (dx, dy) in EMU to add to child positions for slide-absolute coordinates."""
    try:
        elem = group_shape._element
        grpSpPr = elem.find(qn("p:grpSpPr"))
        if grpSpPr is None:
            return (group_shape.left or 0), (group_shape.top or 0)
        xfrm = grpSpPr.find(qn("a:xfrm"))
        if xfrm is None:
            return (group_shape.left or 0), (group_shape.top or 0)
        off = xfrm.find(qn("a:off"))
        ch_off = xfrm.find(qn("a:chOff"))
        off_x = int(off.get("x", 0)) if off is not None else (group_shape.left or 0)
        off_y = int(off.get("y", 0)) if off is not None else (group_shape.top or 0)
        ch_off_x = int(ch_off.get("x", 0)) if ch_off is not None else 0
        ch_off_y = int(ch_off.get("y", 0)) if ch_off is not None else 0
        return off_x - ch_off_x, off_y - ch_off_y
    except Exception:
        return (group_shape.left or 0), (group_shape.top or 0)


# ── Shapes & Images Collection ───────────────────────────────────────────────

def _extract_shapes_and_images(
    shapes, part, prefix: str, theme_colors: Dict[str, str], seen_blobs: set,
    offset_left: int = 0, offset_top: int = 0,
) -> Tuple[List[Dict], List[Dict], List[Dict]]:
    """Extract shapes, images, and text elements from a shape collection."""
    shape_list = []
    image_list = []
    text_list = []

    for shape in shapes:
        try:
            if shape.shape_type == 6 and hasattr(shape, "shapes"):  # GROUP
                dx, dy = _get_group_offset_emu(shape)
                sub_shapes, sub_images, sub_texts = _extract_shapes_and_images(
                    shape.shapes, part, prefix, theme_colors, seen_blobs,
                    offset_left=offset_left + dx, offset_top=offset_top + dy,
                )
                shape_list.extend(sub_shapes)
                image_list.extend(sub_images)
                text_list.extend(sub_texts)
                continue

            pos = {
                "left": _emu_to_inches((shape.left or 0) + offset_left),
                "top": _emu_to_inches((shape.top or 0) + offset_top),
                "width": _emu_to_inches(shape.width),
                "height": _emu_to_inches(shape.height),
            }

            if shape.has_table:
                table_data = _extract_table(shape, theme_colors)
                shape_info: Dict[str, Any] = {
                    "name": shape.name,
                    "type": "TABLE",
                    **pos,
                    "table": table_data,
                }
                shape_list.append(shape_info)
                continue

            is_picture = shape.shape_type == 13
            has_blip = shape._element.find(".//" + qn("a:blip")) is not None
            SVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
            has_svg = shape._element.find(".//{%s}svgBlip" % SVG_NS) is not None

            if is_picture or has_blip or has_svg:
                img_entry = _extract_image_from_shape(shape, part, prefix, seen_blobs,
                                                      offset_left=offset_left, offset_top=offset_top)
                if img_entry is not None:
                    image_list.append(img_entry)
                    if not is_picture and shape.has_text_frame and shape.text_frame.text.strip():
                        pass
                    else:
                        continue

            if shape.has_text_frame and shape.text_frame.text.strip():
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        paragraphs.append(_extract_paragraph(para, theme_colors, ph_shape=shape if shape.is_placeholder else None))

                if paragraphs:
                    te: Dict[str, Any] = {
                        "shape_name": shape.name,
                        **pos,
                        "paragraphs": paragraphs,
                    }
                    try:
                        if shape.placeholder_format is not None:
                            ph_type = str(shape.placeholder_format.type).split(".")[-1].split("(")[0]
                            te["placeholder_type"] = ph_type
                    except Exception:
                        pass
                    try:
                        body_props = _extract_text_body_properties(shape._element)
                        if body_props.get("vertical_align"):
                            te["vertical_align"] = body_props["vertical_align"]
                    except Exception:
                        pass
                    text_list.append(te)
                    continue

            shape_info: Dict[str, Any] = {
                "name": shape.name,
                "type": _shape_type_name(shape),
                **pos,
            }
            try:
                fill_info = _extract_fill(shape, part, theme_colors)
                if fill_info["type"] != "none":
                    shape_info["fill"] = fill_info
            except Exception:
                pass
            try:
                line_props = _extract_line_properties(shape, theme_colors)
                if line_props:
                    shape_info["line"] = line_props
            except Exception:
                pass
            try:
                if shape.rotation and shape.rotation != 0:
                    shape_info["rotation"] = round(shape.rotation, 1)
            except Exception:
                pass

            shape_list.append(shape_info)

        except Exception:
            continue

    return shape_list, image_list, text_list


# ── Placeholder Extraction ───────────────────────────────────────────────────

def _extract_placeholders(placeholders, theme_colors: Dict[str, str]) -> List[Dict]:
    """Extract placeholder definitions from a layout or master."""
    result = []
    for ph in placeholders:
        try:
            ph_info: Dict[str, Any] = {
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type).split(".")[-1].split("(")[0],
                "left": _emu_to_inches(ph.left),
                "top": _emu_to_inches(ph.top),
                "width": _emu_to_inches(ph.width),
                "height": _emu_to_inches(ph.height),
            }

            font_found = False
            if ph.has_text_frame:
                for para in ph.text_frame.paragraphs:
                    if para.runs:
                        fi = _extract_font_info(para.runs[0])
                        if fi:
                            ph_info["font"] = fi
                            font_found = True
                        break
                try:
                    if ph.text_frame.paragraphs and ph.text_frame.paragraphs[0].alignment is not None:
                        ph_info["alignment"] = str(ph.text_frame.paragraphs[0].alignment).split(".")[-1].split("(")[0].lower()
                except Exception:
                    pass

            if not font_found:
                xml_font = _extract_default_font_from_xml(ph._element, theme_colors)
                if xml_font:
                    ph_info["font"] = xml_font
                else:
                    ph_info["font"] = {}

            body_props = _extract_text_body_properties(ph._element)
            if body_props:
                ph_info["text_body"] = body_props

            result.append(ph_info)
        except Exception:
            continue
    return result


# ── Fonts Summary ────────────────────────────────────────────────────────────

def _build_fonts_summary(slides_data: List[Dict], layouts_data: List[Dict]) -> Dict[str, Any]:
    """Aggregate font usage across all slides and layouts."""
    usage: Dict[str, Dict[str, set]] = {}

    def _collect_from_paragraphs(paragraphs: List[Dict], context: str):
        for p in paragraphs:
            font = p.get("font", {})
            name = font.get("name")
            if not name:
                continue
            if name not in usage:
                usage[name] = {"contexts": set(), "sizes_pt": set(), "colors": set()}
            usage[name]["contexts"].add(context)
            if "size_pt" in font:
                usage[name]["sizes_pt"].add(font["size_pt"])
            if "color" in font:
                usage[name]["colors"].add(font["color"])

    for slide in slides_data:
        for te in slide.get("text_elements", []):
            ctx = te.get("placeholder_type", "body").lower()
            _collect_from_paragraphs(te.get("paragraphs", []), ctx)

    for layout in layouts_data:
        for ph in layout.get("placeholders", []):
            if "font" in ph:
                name = ph["font"].get("name")
                if name:
                    if name not in usage:
                        usage[name] = {"contexts": set(), "sizes_pt": set(), "colors": set()}
                    usage[name]["contexts"].add(ph.get("type", "body").lower())
                    if "size_pt" in ph["font"]:
                        usage[name]["sizes_pt"].add(ph["font"]["size_pt"])
                    if "color" in ph["font"]:
                        usage[name]["colors"].add(ph["font"]["color"])

    usage_map = {}
    for font_name, data in usage.items():
        usage_map[font_name] = {
            "contexts": sorted(data["contexts"]),
            "sizes_pt": sorted(data["sizes_pt"]),
            "colors": sorted(data["colors"]),
        }

    return {
        "families_used": sorted(usage.keys()),
        "usage_map": usage_map,
    }


# ── Master Text Styles ───────────────────────────────────────────────────────

def _extract_master_text_styles(master, theme_colors: Dict[str, str]) -> Dict[str, Any]:
    """Extract default text styles from slide master's <p:txStyles> element."""
    styles: Dict[str, Any] = {}
    try:
        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

        tx_styles = master._element.find(f"{{{ns_p}}}txStyles")
        if tx_styles is None:
            return styles

        style_map = {
            f"{{{ns_p}}}titleStyle": "title_style",
            f"{{{ns_p}}}bodyStyle": "body_style",
            f"{{{ns_p}}}otherStyle": "other_style",
        }

        for xml_tag, key in style_map.items():
            style_elem = tx_styles.find(xml_tag)
            if style_elem is None:
                continue

            levels = {}
            for lvl_num in range(1, 10):
                lvl_elem = style_elem.find(f"{{{ns_a}}}lvl{lvl_num}pPr")
                if lvl_elem is None:
                    continue

                lvl_info: Dict[str, Any] = {}

                algn = lvl_elem.get("algn")
                if algn:
                    lvl_info["alignment"] = algn

                indent = lvl_elem.get("indent")
                if indent:
                    lvl_info["indent_pt"] = round(int(indent) / 12700, 2)

                marL = lvl_elem.get("marL")
                if marL:
                    lvl_info["margin_left_pt"] = round(int(marL) / 12700, 2)

                buChar = lvl_elem.find(f"{{{ns_a}}}buChar")
                buAutoNum = lvl_elem.find(f"{{{ns_a}}}buAutoNum")
                buNone = lvl_elem.find(f"{{{ns_a}}}buNone")
                buFont = lvl_elem.find(f"{{{ns_a}}}buFont")

                if buChar is not None:
                    bullet_info: Dict[str, Any] = {"type": "char", "char": buChar.get("char", "")}
                    if buFont is not None:
                        bullet_info["font"] = buFont.get("typeface", "")
                    lvl_info["bullet"] = bullet_info
                elif buAutoNum is not None:
                    lvl_info["bullet"] = {"type": "auto_num", "scheme": buAutoNum.get("type", "")}
                elif buNone is not None:
                    lvl_info["bullet"] = None

                def_rpr = lvl_elem.find(f"{{{ns_a}}}defRPr")
                if def_rpr is not None:
                    font_info: Dict[str, Any] = {}

                    sz = def_rpr.get("sz")
                    if sz:
                        font_info["size_pt"] = round(int(sz) / 100, 1)

                    b = def_rpr.get("b")
                    if b is not None:
                        font_info["bold"] = b == "1"

                    i_attr = def_rpr.get("i")
                    if i_attr is not None:
                        font_info["italic"] = i_attr == "1"

                    latin = def_rpr.find(f"{{{ns_a}}}latin")
                    if latin is not None:
                        typeface = latin.get("typeface", "")
                        if typeface and not typeface.startswith("+"):
                            font_info["name"] = typeface
                        elif typeface == "+mj-lt":
                            font_info["theme_font"] = "major"
                        elif typeface == "+mn-lt":
                            font_info["theme_font"] = "minor"

                    solid_fill = def_rpr.find(f"{{{ns_a}}}solidFill")
                    if solid_fill is not None:
                        srgb = solid_fill.find(f"{{{ns_a}}}srgbClr")
                        if srgb is not None:
                            color_val = srgb.get("val")
                            if color_val:
                                color_val = _apply_color_modifiers(color_val, srgb)
                            font_info["color"] = color_val
                        else:
                            scheme_clr = solid_fill.find(f"{{{ns_a}}}schemeClr")
                            if scheme_clr is not None:
                                clr_name = scheme_clr.get("val", "")
                                resolved = theme_colors.get(clr_name)
                                if resolved:
                                    resolved = _apply_color_modifiers(resolved, scheme_clr)
                                font_info["color"] = resolved
                                font_info["color_theme_ref"] = clr_name

                    if font_info:
                        lvl_info["font"] = font_info

                if lvl_info:
                    levels[f"level{lvl_num}"] = lvl_info

            if levels:
                styles[key] = levels
    except Exception:
        pass
    return styles


# ── Background Extraction ────────────────────────────────────────────────────

def _get_own_bg_element(slide_or_layout):
    """Return the <p:bg> element ONLY if explicitly defined at this level."""
    try:
        cSld = slide_or_layout._element.find(qn("p:cSld"))
        if cSld is not None:
            return cSld.find(qn("p:bg"))
    except Exception:
        pass
    return None


def _extract_background(bg_element, part, theme_colors: Dict[str, str], seen_blobs: set) -> Tuple[Dict[str, Any], List[Dict]]:
    """Extract background info including image fills."""
    images = []
    bg_info: Dict[str, Any] = {"type": "none"}

    try:
        blip_fill = bg_element.find(".//" + qn("a:blipFill"))
        if blip_fill is not None:
            blip = blip_fill.find(qn("a:blip"))
            if blip is not None:
                r_embed = blip.get(qn("r:embed"))
                image_part = _get_image_from_rel(part, r_embed)
                if image_part is not None:
                    blob = image_part.blob
                    blob_hash = hashlib.md5(blob[:512]).hexdigest()
                    ext = image_part.content_type.split("/")[-1].replace("jpeg", "jpg")
                    img_name = f"bg_{blob_hash[:8]}.{ext}"

                    if blob_hash not in seen_blobs:
                        seen_blobs.add(blob_hash)
                        with Image.open(io.BytesIO(blob)) as img:
                            w_px, h_px = img.size
                        images.append({
                            "name": img_name,
                            "media_type": image_part.content_type,
                            "base64_data": base64.b64encode(blob).decode("utf-8"),
                            "width_px": w_px,
                            "height_px": h_px,
                        })

                    bg_info = {"type": "image", "image_ref": img_name}
                    return bg_info, images

        solid_fill = bg_element.find(".//" + qn("a:solidFill"))
        if solid_fill is not None:
            srgb = solid_fill.find(qn("a:srgbClr"))
            if srgb is not None:
                color_val = srgb.get("val")
                if color_val:
                    color_val = _apply_color_modifiers(color_val, srgb)
                bg_info = {"type": "solid", "color": color_val}
                return bg_info, images
            scheme_clr = solid_fill.find(qn("a:schemeClr"))
            if scheme_clr is not None:
                clr_name = scheme_clr.get("val", "")
                resolved = theme_colors.get(clr_name)
                if resolved:
                    resolved = _apply_color_modifiers(resolved, scheme_clr)
                bg_info = {"type": "solid", "color": resolved, "theme_ref": clr_name}
                return bg_info, images

        grad_fill = bg_element.find(".//" + qn("a:gradFill"))
        if grad_fill is not None:
            stops = []
            for gs in grad_fill.findall(".//" + qn("a:gs")):
                pos = gs.get("pos")
                position = int(pos) / 100000 if pos else 0
                color = None
                srgb = gs.find(".//" + qn("a:srgbClr"))
                if srgb is not None:
                    color = srgb.get("val")
                    if color:
                        color = _apply_color_modifiers(color, srgb)
                else:
                    scheme_clr = gs.find(".//" + qn("a:schemeClr"))
                    if scheme_clr is not None:
                        clr_name = scheme_clr.get("val", "")
                        color = theme_colors.get(clr_name)
                        if color:
                            color = _apply_color_modifiers(color, scheme_clr)
                stops.append({"color": color, "position": round(position, 2)})
            bg_info = {"type": "gradient", "stops": stops}
            lin = grad_fill.find(qn("a:lin"))
            if lin is not None:
                ang = lin.get("ang")
                if ang:
                    bg_info["angle"] = round(int(ang) / 60000)
                bg_info["gradient_type"] = "linear"
            else:
                path_el = grad_fill.find(qn("a:path"))
                if path_el is not None:
                    bg_info["gradient_type"] = path_el.get("path", "circle")
                else:
                    bg_info["gradient_type"] = "linear"
            return bg_info, images

        bg_ref = bg_element.find(qn("p:bgRef"))
        if bg_ref is not None:
            scheme_clr = bg_ref.find(qn("a:schemeClr"))
            if scheme_clr is not None:
                clr_name = scheme_clr.get("val", "")
                resolved = theme_colors.get(clr_name)
                if resolved:
                    resolved = _apply_color_modifiers(resolved, scheme_clr)
                bg_info = {"type": "solid", "color": resolved, "theme_ref": clr_name}
                return bg_info, images
            srgb = bg_ref.find(qn("a:srgbClr"))
            if srgb is not None:
                color_val = srgb.get("val")
                if color_val:
                    color_val = _apply_color_modifiers(color_val, srgb)
                bg_info = {"type": "solid", "color": color_val}
                return bg_info, images

    except Exception:
        pass

    return bg_info, images


# ── Effective Background ─────────────────────────────────────────────────────

def _compute_effective_background(
    slide_info: Dict[str, Any],
    layouts_data: List[Dict],
    masters_data: List[Dict],
    slide_width: float = 13.33,
    slide_height: float = 7.5,
) -> Dict[str, Any]:
    """Resolve the visual background by walking the inheritance chain."""
    min_cover_w = slide_width * 0.95
    min_cover_h = slide_height * 0.95

    def _bg_color(bg: Dict[str, Any]) -> Optional[str]:
        if bg.get("type") == "solid" and bg.get("color"):
            return bg["color"]
        return None

    def _covering_color(data: Dict[str, Any]) -> Optional[str]:
        for shape in data.get("shapes", []):
            w = shape.get("width") or 0
            h = shape.get("height") or 0
            fill = shape.get("fill", {})
            if w >= min_cover_w and h >= min_cover_h and fill.get("type") == "solid" and fill.get("color"):
                return fill["color"]
        return None

    layout_idx = slide_info.get("layout_index")
    layout = layouts_data[layout_idx] if layout_idx is not None and layout_idx < len(layouts_data) else None
    master_idx = layout.get("master_index", 0) if layout else 0
    master = masters_data[master_idx] if master_idx < len(masters_data) else None

    slide_cover = _covering_color(slide_info)
    if slide_cover:
        return {"color": slide_cover, "source": "slide_shape"}

    layout_cover = _covering_color(layout) if layout else None
    if layout_cover:
        return {"color": layout_cover, "source": "layout_shape"}

    master_cover = _covering_color(master) if master else None

    slide_bg = slide_info.get("background", {})
    declared_color = _bg_color(slide_bg)
    declared_source = "slide"

    if declared_color is None and layout:
        layout_bg = layout.get("background", {})
        declared_color = _bg_color(layout_bg)
        declared_source = "layout"

    if declared_color is None and master:
        master_bg = master.get("background", {})
        declared_color = _bg_color(master_bg)
        declared_source = "master"

    if master_cover:
        return {
            "color": master_cover,
            "source": "master_shape",
            "declared_color": declared_color,
            "declared_source": declared_source,
        }

    return {
        "color": declared_color or "FFFFFF",
        "source": declared_source if declared_color else "default",
    }


# ── Main Extraction ──────────────────────────────────────────────────────────

def extract_template_context(pptx_path: str, output_dir: str) -> None:
    """Extract template context from a PPTX file and write to output_dir."""
    pptx_path_obj = Path(pptx_path)
    output_path = Path(output_dir)
    images_dir = output_path / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    pptx_bytes = pptx_path_obj.read_bytes()

    try:
        buf = io.BytesIO(pptx_bytes)
        prs = Presentation(buf)
    except Exception as e:
        print(f"ERROR: Failed to parse PPTX: {e}", file=sys.stderr)
        sys.exit(3)

    seen_blobs: set = set()
    all_images: List[Dict] = []

    # ── Presentation info ──
    presentation_info = {
        "slide_width_inches": round(prs.slide_width / 914400, 2),
        "slide_height_inches": round(prs.slide_height / 914400, 2),
        "slide_count": len(prs.slides),
    }

    # ── Theme ──
    themes_per_master = _extract_themes_per_master(pptx_bytes)
    theme = themes_per_master[0] if themes_per_master else _extract_theme(pptx_bytes)
    theme_colors = theme.get("color_scheme", {})

    # ── Slide Masters ──
    masters_data = []
    resolved_colors_per_master: List[Dict[str, str]] = []
    for m_idx, master in enumerate(prs.slide_masters):
        master_theme = themes_per_master[m_idx] if m_idx < len(themes_per_master) else theme
        master_theme_colors = master_theme.get("color_scheme", {})
        clr_map = _extract_clr_map(master)
        resolved_tc = _build_resolved_theme_colors(master_theme_colors, clr_map)
        resolved_colors_per_master.append(resolved_tc)

        m_info: Dict[str, Any] = {"index": m_idx}
        try:
            m_info["name"] = master.name if hasattr(master, "name") else f"Master {m_idx}"
        except Exception:
            m_info["name"] = f"Master {m_idx}"

        own_bg = _get_own_bg_element(master)
        if own_bg is not None:
            try:
                bg_info, bg_images = _extract_background(own_bg, master.part, resolved_tc, seen_blobs)
                m_info["background"] = bg_info
                all_images.extend(bg_images)
            except Exception:
                m_info["background"] = {"type": "none"}
        else:
            m_info["background"] = {"type": "none"}

        shapes, images, texts = _extract_shapes_and_images(
            master.shapes, master.part, f"master{m_idx}", resolved_tc, seen_blobs
        )
        m_info["shapes"] = shapes
        m_info["images"] = [{k: v for k, v in img.items() if k != "base64_data"} for img in images]
        all_images.extend(images)

        text_styles = _extract_master_text_styles(master, resolved_tc)
        if text_styles:
            m_info["text_styles"] = text_styles

        masters_data.append(m_info)

    # ── Slide Layouts ──
    layouts_data = []
    layout_idx = 0
    for m_idx, master in enumerate(prs.slide_masters):
        resolved_tc = resolved_colors_per_master[m_idx] if m_idx < len(resolved_colors_per_master) else theme_colors
        for layout in master.slide_layouts:
            l_info: Dict[str, Any] = {
                "index": layout_idx,
                "name": layout.name,
                "master_index": m_idx,
            }

            own_bg = _get_own_bg_element(layout)
            if own_bg is not None:
                try:
                    bg_info, bg_images = _extract_background(own_bg, layout.part, resolved_tc, seen_blobs)
                    l_info["background"] = bg_info
                    all_images.extend(bg_images)
                except Exception:
                    l_info["background"] = {"type": "inherit"}
            else:
                l_info["background"] = {"type": "inherit"}

            l_info["placeholders"] = _extract_placeholders(layout.placeholders, resolved_tc)

            shapes, images, texts = _extract_shapes_and_images(
                layout.shapes, layout.part, f"layout{layout_idx}", resolved_tc, seen_blobs
            )
            l_info["shapes"] = shapes
            l_info["images"] = [{k: v for k, v in img.items() if k != "base64_data"} for img in images]
            all_images.extend(images)

            layouts_data.append(l_info)
            layout_idx += 1

    # ── Slides ──
    slides_data = []
    for s_idx, slide in enumerate(prs.slides, start=1):
        s_info: Dict[str, Any] = {"index": s_idx}

        try:
            layout_name = slide.slide_layout.name
            s_info["layout_name"] = layout_name
            for li, ld in enumerate(layouts_data):
                if ld["name"] == layout_name:
                    s_info["layout_index"] = li
                    break
        except Exception:
            s_info["layout_name"] = "Unknown"

        slide_master_idx = 0
        li = s_info.get("layout_index")
        if li is not None and li < len(layouts_data):
            slide_master_idx = layouts_data[li].get("master_index", 0)
        slide_resolved_tc = resolved_colors_per_master[slide_master_idx] if slide_master_idx < len(resolved_colors_per_master) else theme_colors

        own_bg = _get_own_bg_element(slide)
        if own_bg is not None:
            try:
                bg_info, bg_images = _extract_background(own_bg, slide.part, slide_resolved_tc, seen_blobs)
                s_info["background"] = bg_info
                all_images.extend(bg_images)
            except Exception:
                s_info["background"] = {"type": "inherit"}
        else:
            s_info["background"] = {"type": "inherit"}

        shapes, images, texts = _extract_shapes_and_images(
            slide.shapes, slide.part, f"slide{s_idx}", slide_resolved_tc, seen_blobs
        )
        s_info["shapes"] = shapes
        s_info["images"] = [{k: v for k, v in img.items() if k != "base64_data"} for img in images]
        s_info["text_elements"] = texts
        all_images.extend(images)

        slides_data.append(s_info)

    # ── Effective background per slide ──
    for s_info in slides_data:
        eff_bg = _compute_effective_background(
            s_info, layouts_data, masters_data,
            slide_width=presentation_info["slide_width_inches"],
            slide_height=presentation_info["slide_height_inches"],
        )
        s_info["effective_background"] = eff_bg

    # ── Fonts summary ──
    fonts_summary = _build_fonts_summary(slides_data, layouts_data)

    # ── Images manifest ──
    images_manifest = []
    for img in all_images:
        entry: Dict[str, Any] = {"ref": img["name"]}
        if "media_type" in img:
            entry["media_type"] = img["media_type"]
        if "width_px" in img:
            entry["width_px"] = img["width_px"]
            entry["height_px"] = img["height_px"]
        name = img["name"]
        if name.startswith("master"):
            entry["source"] = "master"
        elif name.startswith("layout"):
            entry["source"] = "layout"
        elif name.startswith("slide"):
            entry["source"] = "slide"
        elif name.startswith("bg_"):
            entry["source"] = "background"
        images_manifest.append(entry)

    # ── Build context ──
    context = {
        "presentation": presentation_info,
        "theme": theme,
        "slide_masters": masters_data,
        "slide_layouts": layouts_data,
        "slides": slides_data,
        "fonts_summary": fonts_summary,
        "images_manifest": images_manifest,
    }

    # ── Write context.json ──
    context_file = output_path / "context.json"
    context_file.write_text(json.dumps(context, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"Wrote context.json ({context_file.stat().st_size} bytes)")

    # ── Write images ──
    images_with_data = [img for img in all_images if "base64_data" in img]
    written_names: set = set()
    for img in images_with_data:
        img_name = img["name"]
        if img_name in written_names:
            continue
        written_names.add(img_name)
        img_path = images_dir / img_name
        img_path.write_bytes(base64.b64decode(img["base64_data"]))

    print(f"Wrote {len(written_names)} image(s) to images/")
    print(f"Done. Masters: {len(masters_data)}, Layouts: {len(layouts_data)}, Slides: {len(slides_data)}")


# ── CLI Entry Point ──────────────────────────────────────────────────────────

def main():
    if len(sys.argv) != 3:
        print("Usage: python extract_template.py <pptx_path> <output_dir>", file=sys.stderr)
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_dir = sys.argv[2]

    if not os.path.isfile(pptx_path):
        print(f"ERROR: File not found: {pptx_path}", file=sys.stderr)
        sys.exit(2)

    if not pptx_path.lower().endswith(".pptx"):
        print(f"ERROR: File must be a .pptx file: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    extract_template_context(pptx_path, output_dir)


if __name__ == "__main__":
    main()
